#!/usr/bin/env python3
"""
案例PPT生成脚本
基于文心案例PPT模板生成标准化的案例PPT
"""

import os
import sys
import json
from typing import Dict, List, Optional, Any
from dataclasses import dataclass


@dataclass
class CasePPTData:
    """案例PPT数据结构"""
    title: str  # 案例标题
    case_description: str  # 案例说明文字（250-350字）
    model_label: str = "文心大模型应用"  # 模型标签
    effect_metrics: List[Dict[str, str]] = None  # 效果指标列表
    
    def __post_init__(self):
        if self.effect_metrics is None:
            self.effect_metrics = []


class CasePPTGenerator:
    """案例PPT生成器"""
    
    def __init__(self, template_path: Optional[str] = None):
        """
        初始化PPT生成器
        
        Args:
            template_path: PPT模板文件路径，如果为None则需要手动指定
        """
        self.template_path = template_path
        self.required_elements = {
            'title': '案例标题',
            'case_description': '案例说明',
            'model_label': '模型标签',
            'effect_metrics': '效果指标'
        }
    
    def validate_case_data(self, case_data: CasePPTData) -> Dict[str, Any]:
        """验证案例数据完整性"""
        validation_result = {
            'valid': True,
            'missing': [],
            'warnings': []
        }
        
        # 检查必填字段
        if not case_data.title:
            validation_result['missing'].append('案例标题')
            validation_result['valid'] = False
        
        if not case_data.case_description:
            validation_result['missing'].append('案例说明')
            validation_result['valid'] = False
        
        # 检查案例说明字数
        desc_length = len(case_data.case_description)
        if desc_length < 250:
            validation_result['warnings'].append(f'案例说明字数不足250字（当前{desc_length}字）')
        elif desc_length > 350:
            validation_result['warnings'].append(f'案例说明字数超过350字（当前{desc_length}字）')
        
        # 检查效果指标
        if not case_data.effect_metrics:
            validation_result['warnings'].append('未提供效果指标，建议添加量化数据')
        
        return validation_result
    
    def extract_effect_metrics_from_text(self, text: str) -> List[Dict[str, str]]:
        """
        从案例文字中提取效果指标
        
        Args:
            text: 案例说明文字
        
        Returns:
            效果指标列表，格式：[{'value': '95%', 'label': '效率提升'}]
        """
        import re
        
        metrics = []
        
        # 提取百分比提升
        percentage_pattern = r'(\d+(?:\.\d+)?)%\s*[提升降低减少增加]'
        percentage_matches = re.findall(percentage_pattern, text)
        
        # 提取时间对比
        time_pattern = r'(\d+(?:\.\d+)?)\s*(小时|分钟|天|月)\s*[降至→→降为]?\s*(\d+(?:\.\d+)?)\s*(小时|分钟|天|月)'
        time_matches = re.findall(time_pattern, text)
        
        # 提取具体数字
        number_pattern = r'(\d+(?:\.\d+)?)\s*(个|项|次|套|用户|场景)'
        number_matches = re.findall(number_pattern, text)
        
        # 生成指标（示例）
        if percentage_matches:
            for i, match in enumerate(percentage_matches[:2]):  # 最多2个
                metrics.append({
                    'value': f'{match}%',
                    'label': '效率提升' if i == 0 else '准确率'
                })
        
        if time_matches and len(metrics) < 3:
            before, unit1, after, unit2 = time_matches[0]
            metrics.append({
                'value': f'{after}{unit2}',
                'label': '处理时间'
            })
        
        if number_matches and len(metrics) < 3:
            number, unit = number_matches[0]
            metrics.append({
                'value': f'{number}{unit}',
                'label': '项目规模'
            })
        
        return metrics
    
    def generate_ppt_content(self, case_data: CasePPTData) -> Dict[str, Any]:
        """
        生成PPT内容结构
        
        Args:
            case_data: 案例数据
        
        Returns:
            PPT内容结构
        """
        # 验证数据
        validation = self.validate_case_data(case_data)
        if not validation['valid']:
            return {
                'success': False,
                'error': f"数据验证失败，缺少必填项: {', '.join(validation['missing'])}",
                'warnings': validation['warnings']
            }
        
        # 如果没有提供效果指标，尝试从文字中提取
        if not case_data.effect_metrics:
            case_data.effect_metrics = self.extract_effect_metrics_from_text(case_data.case_description)
        
        # 构建PPT内容结构
        ppt_content = {
            'slide_layout': 'case_template',  # 使用案例模板布局
            'elements': {
                'title': {
                    'text': case_data.title,
                    'position': 'top-center',
                    'font_size': 32,
                    'font_weight': 'bold'
                },
                'logo': {
                    'position': 'top-right',
                    'type': 'image',
                    'required': True
                },
                'case_description': {
                    'text': case_data.case_description,
                    'position': 'left-body',
                    'font_size': 14,
                    'line_height': 1.5
                },
                'model_label': {
                    'text': case_data.model_label,
                    'position': 'right-top',
                    'background_color': '#4CAF50',  # 绿色
                    'font_color': '#FFFFFF',
                    'font_size': 12,
                    'padding': '8px 16px'
                },
                'effect_metrics': {
                    'position': 'right-body',
                    'metrics': case_data.effect_metrics,
                    'layout': 'vertical',  # 垂直排列
                    'style': {
                        'value_font_size': 36,
                        'value_font_weight': 'bold',
                        'value_color': '#1E88E5',  # 蓝色
                        'label_font_size': 14,
                        'label_color': '#666666'
                    }
                }
            },
            'color_scheme': {
                'primary': '#1E88E5',  # 主色调：蓝色
                'secondary': '#4CAF50',  # 辅助色：绿色
                'background': '#FFFFFF',  # 背景：白色
                'text': '#333333'  # 文字：深灰色
            }
        }
        
        return {
            'success': True,
            'ppt_content': ppt_content,
            'warnings': validation['warnings'],
            'message': 'PPT内容生成成功'
        }
    
    def create_ppt_file(self, ppt_content: Dict[str, Any], output_path: str) -> Dict[str, Any]:
        """
        创建PPT文件
        
        注意：此方法需要python-pptx库支持
        实际使用时需要安装：pip install python-pptx
        
        Args:
            ppt_content: PPT内容结构
            output_path: 输出文件路径
        
        Returns:
            生成结果
        """
        try:
            # 检查是否安装了python-pptx库
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor
            except ImportError:
                return {
                    'success': False,
                    'error': '缺少python-pptx库，请运行: pip install python-pptx',
                    'alternative': '可以使用生成的内容结构，手动在PowerPoint中创建PPT'
                }
            
            # 创建PPT
            prs = Presentation()
            
            # 添加空白幻灯片
            blank_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(blank_layout)
            
            # 提取内容
            elements = ppt_content['elements']
            
            # 添加标题
            # 实际实现需要根据模板的具体布局调整位置和大小
            # 这里提供一个简化的示例
            
            # 保存PPT
            prs.save(output_path)
            
            return {
                'success': True,
                'output_path': output_path,
                'message': f'PPT文件已生成: {output_path}'
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'PPT生成失败: {str(e)}'
            }
    
    def generate_ppt_from_case_text(self, case_text: str, title: str, output_path: str) -> Dict[str, Any]:
        """
        从案例文字直接生成PPT
        
        Args:
            case_text: 案例说明文字
            title: 案例标题
            output_path: 输出文件路径
        
        Returns:
            生成结果
        """
        # 创建案例数据
        case_data = CasePPTData(
            title=title,
            case_description=case_text
        )
        
        # 提取效果指标
        case_data.effect_metrics = self.extract_effect_metrics_from_text(case_text)
        
        # 生成PPT内容
        content_result = self.generate_ppt_content(case_data)
        
        if not content_result['success']:
            return content_result
        
        # 创建PPT文件
        file_result = self.create_ppt_file(content_result['ppt_content'], output_path)
        
        return {
            **file_result,
            'warnings': content_result.get('warnings', [])
        }


def main():
    """命令行入口"""
    import argparse
    
    parser = argparse.ArgumentParser(description='文心案例PPT生成工具')
    parser.add_argument('--title', required=True, help='案例标题')
    parser.add_argument('--text', required=True, help='案例说明文字（250-350字）')
    parser.add_argument('--output', default='./output.pptx', help='输出文件路径')
    parser.add_argument('--template', help='PPT模板路径（可选）')
    
    args = parser.parse_args()
    
    # 创建生成器
    generator = CasePPTGenerator(template_path=args.template)
    
    # 生成PPT
    result = generator.generate_ppt_from_case_text(
        case_text=args.text,
        title=args.title,
        output_path=args.output
    )
    
    if result['success']:
        print(f"✅ {result['message']}")
        if result.get('warnings'):
            print("⚠️  警告:")
            for warning in result['warnings']:
                print(f"   - {warning}")
    else:
        print(f"❌ {result['error']}")
        if result.get('alternative'):
            print(f"💡 替代方案: {result['alternative']}")


if __name__ == "__main__":
    main()